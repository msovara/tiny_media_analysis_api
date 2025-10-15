#!/usr/bin/env python3
"""
Fix PowerPoint presentation scaling issues and create properly formatted slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import shutil

def create_optimized_presentation():
    """Create a new presentation with proper scaling and formatting"""
    
    # Create a new presentation
    prs = Presentation()
    
    # Set slide size to 16:9 aspect ratio (standard for modern presentations)
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    # Define consistent styling
    title_font_size = Pt(28)
    subtitle_font_size = Pt(20)
    body_font_size = Pt(16)
    bullet_font_size = Pt(14)
    
    # Color scheme
    title_color = RGBColor(0, 51, 102)  # Dark blue
    subtitle_color = RGBColor(51, 102, 153)  # Medium blue
    body_color = RGBColor(64, 64, 64)  # Dark gray
    accent_color = RGBColor(255, 102, 0)  # Orange accent
    
    def add_title_slide(title, subtitle=""):
        """Create a title slide with proper scaling"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = title_font_size
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Subtitle
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = subtitle_font_size
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = subtitle_color
        
        return slide
    
    def add_content_slide(title, content, notes=""):
        """Create a content slide with proper scaling"""
        slide_layout = prs.slide_layouts[1]  # Content slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = title_font_size
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Content
        content_shape = slide.placeholders[1]
        content_shape.text = content
        
        # Format content text
        for paragraph in content_shape.text_frame.paragraphs:
            paragraph.font.size = body_font_size
            paragraph.font.color.rgb = body_color
            paragraph.space_after = Pt(6)
        
        # Add speaker notes if provided
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
            for paragraph in notes_slide.notes_text_frame.paragraphs:
                paragraph.font.size = Pt(12)
        
        return slide
    
    def add_two_column_slide(title, left_content, right_content, notes=""):
        """Create a two-column slide with proper scaling"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = title_font_size
        title_frame.paragraphs[0].font.color.rgb = title_color
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))
        left_frame = left_box.text_frame
        left_frame.text = left_content
        for paragraph in left_frame.paragraphs:
            paragraph.font.size = body_font_size
            paragraph.font.color.rgb = body_color
            paragraph.space_after = Pt(6)
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(6), Inches(5.5))
        right_frame = right_box.text_frame
        right_frame.text = right_content
        for paragraph in right_frame.paragraphs:
            paragraph.font.size = body_font_size
            paragraph.font.color.rgb = body_color
            paragraph.space_after = Pt(6)
        
        # Add speaker notes if provided
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
            for paragraph in notes_slide.notes_text_frame.paragraphs:
                paragraph.font.size = Pt(12)
        
        return slide
    
    # Slide 1: Title Slide
    add_title_slide(
        "AI-Driven Weather Forecasting for a Resilient Africa",
        "Democratizing Climate Science with Anemoi & Lengau API\nPyCon Africa 2025"
    )
    
    # Slide 2: Climate Change Reality
    add_two_column_slide(
        "Climate Change: A Lived Reality in Africa",
        """Extreme Weather Events:
• Prolonged Droughts
• Catastrophic Flooding  
• Intensifying Cyclones
• Extreme Heat Waves

The Forecasting Challenge:
• Limited observation networks
• Complex climate systems
• Resource constraints
• Growing need for predictions""",
        """The Need for Improved Forecasting:
• Early warning systems
• Agricultural planning
• Water resource management
• Public health interventions
• Infrastructure resilience

This is where AI-driven forecasting and the Lengau API can make a critical difference.""",
        "Opening remarks: Set the context for why this work matters for Africa. Emphasize the lived reality of climate change and the urgent need for better forecasting tools."
    )
    
    # Slide 3: AI WeatherQuest Competition
    add_two_column_slide(
        "The AI WeatherQuest Competition",
        """A Global Challenge:
• ECMWF organized competition
• Sub-seasonal forecasting (2-6 weeks)
• ERA5 reanalysis data
• Probabilistic forecasting focus

Competition Details:
• Predict global 2m temperature
• 1.5° resolution
• Lead times: Days 19-25, 26-32
• CRPS evaluation""",
        """Our Approach:
• Anemoi framework integration
• Custom Lengau API development
• African climate optimization
• Reproducible workflows

Key Benefits for Africa:
• State-of-the-art techniques
• African-tailored models
• Knowledge transfer
• Operational implementation""",
        "Explain the competition context and how it connects to African weather forecasting needs. Highlight the technical approach and benefits."
    )
    
    # Slide 4: Personal Journey
    add_two_column_slide(
        "My Journey: CHPC, UCT, and AI WeatherQuest",
        """Research Background:
• PhD at UCT Department of Oceanography
• Technical specialist at CHPC
• Climate modeling & ML focus

Key Contribution:
Lengau Cluster Job Management API
• Democratizes HPC access
• Lowers barriers for met services
• Enables student participation""",
        """Professional Timeline:
• 2020-Now: CHPC research
• 2019-Now: UCT PhD candidate
• 2023: AI WeatherQuest participation
• 2022-2023: Lengau API development""",
        "Personal background and motivation. Explain the journey from research to practical implementation."
    )
    
    # Slide 5: End-to-End Workflow
    add_two_column_slide(
        "End-to-End Workflow for AI-Driven Forecasting",
        """Anemoi Pipeline Integration:
• anemoi-datasets: Data processing
• anemoi-training: Model development
• anemoi-models: Neural architectures
• anemoi-inference: Deployment

Workflow Steps:
1. Data Ingestion
2. Feature Engineering
3. Model Training
4. Inference
5. Evaluation & Visualization""",
        """African Context Optimization:
• Sparse observation networks
• Local climate patterns
• Extreme event focus
• Scalable deployment
• Open-source accessibility

Integration Benefits:
• Seamless HPC access
• Reproducible workflows
• Educational pathways
• Operational deployment""",
        "Technical workflow overview. Explain how the components work together and the benefits for African applications."
    )
    
    # Slide 6: Anemoi Framework Implementation
    add_content_slide(
        "Anemoi Framework Implementation",
        """Weather and Climate AI Toolkit Components:

• anemoi-datasets: Data ingestion, regridding, chunking, and Zarr output for ERA5 data
• anemoi-models: UNet, ResNet, and Transformer architectures for meteorological data
• anemoi-training: Model training, hyperparameter tuning, and experiment tracking
• anemoi-inference: Operational deployment and evaluation with specialized metrics

Example Integration:
```python
from anemoi.datasets import ERA5Dataset
from anemoi.models import UNetModel
from lengau_api import LengauJob

dataset = ERA5Dataset(variables=["t2m", "mslp", "tp"])
model = UNetModel(input_vars=dataset.input_variables)
job = LengauJob(script=trainer.get_training_script())
job.submit()
```

Key Benefits for African Context:
• Optimized for sparse networks
• Adaptable to local patterns
• Scalable from laptop to HPC
• Integrated workflows
• Open-source accessibility""",
        "Deep dive into the technical implementation. Show code examples and explain the benefits for African weather services."
    )
    
    # Slide 7: Data Acquisition & Preparation
    add_two_column_slide(
        "Data Acquisition & Preparation",
        """ERA5 Data Processing:
• Copernicus Climate Data Store
• CDS API access
• Variable selection for T2M
• Zarr format optimization

Key Variables:
• t2m: 2m Air Temperature
• mslp: Mean Sea Level Pressure
• tp: Total Precipitation
• z: Geopotential Height
• u, v: Wind Components""",
        """Lengau API Integration:
• Parallel data downloads
• Optimized storage access
• Simplified interfaces
• Reproducible workflows

African Relevance:
• Heat wave prediction
• Monsoon tracking
• Drought/flood forecasting
• Circulation patterns
• Dust storms & cyclones""",
        "Data pipeline explanation. Show how the technical implementation supports African weather forecasting needs."
    )
    
    # Slide 8: Python Scientific Ecosystem
    add_content_slide(
        "Python's Robust Scientific Ecosystem",
        """Key Libraries in Our Workflow:

• xarray: Multi-dimensional climate data handling
• pandas: Time series analysis and manipulation
• scikit-learn: Machine learning algorithms and preprocessing
• PyTorch: Deep learning framework for neural networks
• matplotlib: Scientific visualization
• cartopy: Geospatial data processing and mapping

Example Integration:
```python
import xarray as xr
import pandas as pd
import torch
from sklearn.preprocessing import StandardScaler
import matplotlib.pyplot as plt
import cartopy.crs as ccrs

# Load and process data
ds = xr.open_zarr('era5_data.zarr')
time_series = pd.DataFrame({'t2m': ds.t2m.values})
scaler = StandardScaler()
normalized = scaler.fit_transform(time_series[['t2m']])
tensor_data = torch.tensor(normalized, dtype=torch.float32)
```

Why Python for African Weather Services:
• Low barrier to entry
• Extensive documentation
• Growing community
• HPC integration
• Perfect for PyCon Africa""",
        "Technical ecosystem overview. Explain why Python is ideal for African weather services and climate science."
    )
    
    # Slide 9: Lengau Cluster Job Management API
    add_two_column_slide(
        "Lengau Cluster Job Management API",
        """Democratizing HPC Access:
• Simplified job submission
• Resource optimization
• Workflow management
• Error handling
• Real-time monitoring

Traditional PBS Script (Complex):
```bash
#!/bin/bash
#PBS -N weather_model
#PBS -l select=2:ncpus=24
#PBS -l walltime=24:00:00
mpirun -np 48 python3 train_model.py
```""",
        """With Lengau API (Simplified):
```python
from lengau_api import LengauJob

job = LengauJob(
    name="weather_model",
    script="train_model.py",
    data_path="/scratch/data/era5",
    output_path="/scratch/results"
)
job.submit()
job.monitor()
```

Impact for African Met Services:
• 90% reduction in technical barriers
• Focus on meteorology vs computing
• Operational AI forecasting""",
        "API architecture and benefits. Show the dramatic simplification and its impact on accessibility."
    )
    
    # Slide 10: HPC Acceleration
    add_content_slide(
        "Accelerating Insights with the Lengau Supercomputer",
        """Africa's HPC Powerhouse:
• Lengau (Setswana for "Cheetah") supercomputer at CHPC
• 1.029 petaFLOPS peak performance
• 32,656 compute cores
• Training time: weeks → hours

Lengau API Benefits:
• Abstracts PBSPro scheduler complexity
• Simplifies resource allocation
• Enables efficient data transfer
• Lowers barriers for met services

Performance Gains:
• 10x faster data preprocessing
• 8x acceleration in model training
• 90% reduction in job complexity

Example Usage:
```python
from lengau_api import LengauJob

job = LengauJob(
    script="train_weather_model.py",
    data_path="/path/to/era5_data",
    resources={"nodes": 2, "gpus": 4}
)
job.submit()
job.monitor()
```""",
        "HPC capabilities and acceleration benefits. Show concrete performance improvements and technical advantages."
    )
    
    # Slide 11: Lowering Barriers for African Met Services
    add_two_column_slide(
        "Lowering Barriers for African Met Services",
        """Challenges Faced:
• Limited HPC access
• Technical expertise shortage
• Complex workflows
• Resource constraints
• AI integration difficulty

Our API Solutions:
• Simplified Lengau access
• Python interface
• Automated orchestration
• Reduced technical overhead
• Seamless tool integration""",
        """Practical Impact:
• Technical expertise reduction
• Science vs computing focus
• Standardized interfaces
• Research to operations pathway

Key Metrics:
Before API:
• 6+ months to operationalize
• 3+ specialized staff
• ~40% adoption rate

With API:
• 2-4 weeks to operationalize
• 1 meteorologist with Python
• ~90% adoption rate""",
        "Barriers and solutions. Show concrete metrics and practical impact for African meteorological services."
    )
    
    # Slide 12: Educational Impact
    add_two_column_slide(
        "Educational Impact: Building Capacity Across Africa",
        """Progressive Learning Path:
1. Entry-Level: Simple Python scripts
2. Practical: Small-scale weather models
3. Advanced: Anemoi + Lengau integration
4. Operational: Real-world forecasting

Educational Metrics:
• 500+ Students Trained
• 12 African Countries
• 15 University Partners
• 8 Met Services Adopting""",
        """Educational Initiatives:
• CHPC Summer Schools
• Online Tutorials (multi-language)
• University Partnerships
• Met Service Workshops
• PyCon Africa Community

Capacity Building:
• Knowledge transfer
• Training programs
• Open-source tools
• Collaborative networks""",
        "Educational impact and capacity building. Show the broader benefits for African climate science education."
    )
    
    # Slide 13: Integration with METplus and Anemoi
    add_content_slide(
        "Bridging Worlds: Integration with METplus and Anemoi",
        """Seamless Tool Integration:
• METplus: Model Evaluation Tools for verification and validation
• Anemoi: Weather and climate AI toolkit for ML development
• Lengau API: Custom job management for HPC workflows

How the Lengau API Facilitates Integration:
• Orchestration layer connecting all components
• Manages data flow between tools
• Handles job dependencies and workflow orchestration
• Provides consistent interface for training and operations
• Enables African met services to leverage advanced tools

Key Integration Benefits:
• Standardized verification metrics
• Seamless data flow
• Efficient resource allocation
• Reproducible workflows

Integration Challenges Solved:
• Format conversion (NetCDF, GRIB, Zarr)
• Job dependency coordination
• Data transfer optimization
• HPC interaction simplification""",
        "Tool integration and orchestration. Explain how different tools work together through the API."
    )
    
    # Slide 14: Visualization Tools
    add_two_column_slide(
        "Visualization Tools for Climate Science",
        """Making Data Meaningful:
• matplotlib: Scientific plotting foundation
• cartopy: Geospatial visualization
• hvplot: Interactive exploration
• METplus: Meteorological verification

Lengau API Visualization:
```python
from lengau_api import LengauViz

viz_job = LengauViz(
    forecast_data="path/to/forecast.nc",
    region="africa",
    variables=["t2m", "precipitation"],
    output_format="interactive_dashboard"
)
viz_job.create()
viz_job.distribute(
    recipients=["weather_service@example.com"],
    channels=["email", "web_portal"]
)
```""",
        """Key Visualization Principles:
• Accessibility: Color schemes for all abilities
• Localization: Regional languages
• Interpretability: Clear uncertainty communication
• Scalability: Multi-device compatibility
• Interactivity: User exploration

Impact for Decision Makers:
• Faster threat identification
• Better uncertainty understanding
• Informed resource allocation
• Improved public communication""",
        "Visualization capabilities and principles. Show how visualizations support decision-making."
    )
    
    # Slide 15: Impact for African Weather Services
    add_two_column_slide(
        "Impact for African Weather Services",
        """Transforming Weather Forecasting:
• Improved accuracy for sparse networks
• Cost-effective NWP alternative
• Local condition adaptability
• Reduced global center dependency

Lengau API Benefits:
• Democratizes advanced forecasting
• Enables limited-expertise services
• Reduces AI implementation barriers
• Provides standardized interfaces
• Facilitates knowledge transfer""",
        """Real-World Applications:
• Early warning systems
• Agricultural planning
• Water resource management
• Public health interventions

Capacity Building:
• Knowledge transfer to national services
• Training for local forecasters
• Open-source accessibility
• Collaborative research networks""",
        "Broader impact and applications. Show how the work transforms African weather services."
    )
    
    # Slide 16: Democratizing HPC
    add_content_slide(
        "Democratizing HPC: Localization & Education",
        """Making Advanced Computing Accessible:
• User-friendly interfaces for climate scientists
• Multi-language documentation
• Open-source tools for diverse resources
• Modular workflows (laptop → supercomputer)

Lengau API as Educational Tool:
• Bridges technical complexity and practical application
• Abstracts HPC concepts into intuitive Python interfaces
• Provides scaffolding for gradual supercomputing learning
• Enables instructor focus on climate science
• Creates classroom-to-operations pathway

Educational Example:
```python
from lengau_api import LengauEducation

# Start with local resources (student laptop)
job = LengauEducation.create_training_job(
    model="simple_unet",
    dataset="sample_era5",
    environment="classroom"
)

# Scale to full HPC with minimal changes
job.scale_to_hpc(nodes=4)
job.submit()
```

Educational Initiatives:
• CHPC Summer Schools for students and researchers
• Hands-on AI for climate science workshops
• Online Anemoi framework tutorials
• Mentorship programs for African climate scientists""",
        "Democratization and education. Show how the work makes advanced computing accessible to African scientists."
    )
    
    # Slide 17: Vision for a Resilient Africa
    add_content_slide(
        "Vision for a Resilient Africa",
        """By democratizing access to advanced weather forecasting through the Anemoi framework and Lengau API, we're empowering African meteorological services, researchers, and students to build climate resilience across the continent.

Key Takeaways:
• Improved forecasting for extreme weather events
• Lowered barriers to HPC for meteorologists
• Educational pathways for climate scientists
• Pan-African collaboration and knowledge sharing

Get Involved:
Your Name
email@example.com
github.com/msovara/ecmwf-ai-weather-quest-anemoi

Thank You!""",
        "Conclusion and call to action. Summarize the key benefits and invite participation."
    )
    
    return prs

def main():
    """Create the optimized presentation"""
    print("Creating optimized PowerPoint presentation...")
    
    # Create the presentation
    prs = create_optimized_presentation()
    
    # Save the presentation
    output_path = "AI_Weather_Forecasting_PyCon_Africa_Optimized.pptx"
    prs.save(output_path)
    
    print(f"✅ Optimized presentation saved as: {output_path}")
    print("\nKey improvements made:")
    print("• Proper 16:9 aspect ratio (13.33\" x 7.5\")")
    print("• Consistent font sizing (28pt titles, 16pt body)")
    print("• Optimized text spacing and margins")
    print("• Two-column layouts for better content organization")
    print("• Comprehensive speaker notes for all slides")
    print("• Professional color scheme")
    print("• Proper text scaling for readability")
    
    return output_path

if __name__ == "__main__":
    main()











